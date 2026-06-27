from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("outputs/新能源电力交易决策辅助系统演示.pptx")


def _set_run(run, size=22, bold=False, color=(31, 41, 51)):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = text
    _set_run(p.runs[0], size=26, bold=True, color=(16, 42, 67))


def _bullets(slide, items: list[str], top: float = 1.35):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(5.5))
    frame = box.text_frame
    frame.clear()
    for item in items:
        p = frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        _set_run(p.runs[0], size=20, color=(52, 64, 84))


def _footer(slide):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = "模拟数据演示，不代表真实交易收益，不构成交易建议"
    p.alignment = PP_ALIGN.RIGHT
    _set_run(p.runs[0], size=9, color=(102, 112, 133))


def _add_slide(prs, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(248, 250, 252)
    _title(slide, title)
    _bullets(slide, bullets)
    _footer(slide)
    return slide


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor(16, 42, 67)
    title = cover.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(11.8), Inches(1.0))
    p = title.text_frame.paragraphs[0]
    p.text = "新能源电力交易决策辅助系统"
    _set_run(p.runs[0], size=38, bold=True, color=(255, 255, 255))
    sub = cover.shapes.add_textbox(Inches(0.85), Inches(2.85), Inches(10.8), Inches(1.2))
    p = sub.text_frame.paragraphs[0]
    p.text = "面向集中式光伏电站的日前交易、偏差风险与储能增益演示"
    _set_run(p.runs[0], size=22, color=(217, 226, 236))
    note = cover.shapes.add_textbox(Inches(0.85), Inches(5.75), Inches(10.8), Inches(0.5))
    p = note.text_frame.paragraphs[0]
    p.text = "求职演示作品｜决策辅助，不做自动下单"
    _set_run(p.runs[0], size=16, color=(173, 181, 189))

    slides = [
        (
            "1. 行业变化：光伏从建设逻辑进入交易逻辑",
            [
                "新能源装机增加后，收益不再只取决于发电量和装机规模。",
                "午间光伏出力高时，价格可能承压；傍晚价格高时，光伏自然出力下降。",
                "交易负责人需要同时管理预测、申报、偏差、复盘和储能配置。",
            ],
        ),
        (
            "2. 核心痛点：收益机会与偏差风险并存",
            [
                "出力预测不准，会影响日前申报和结算收益。",
                "只追求高申报量，可能放大偏差考核损失。",
                "只看单日收益，容易忽略波动、极端天气和价格异常。",
                "储能是否值得配，需要基于价差和约束量化判断。",
            ],
        ),
        (
            "3. 系统定位：辅助交易负责人做决策",
            [
                "系统不是自动交易机器人，也不承诺收益。",
                "它把功率预测、电价预测、申报建议、收益回测和储能模拟放在一个闭环里。",
                "目标是让交易判断更系统、更可解释、更便于复盘。",
            ],
        ),
        (
            "4. 业务流程：从预测到申报再到复盘",
            [
                "第一步：预测明日 24 小时光伏出力。",
                "第二步：预测明日 24 小时市场价格和风险时段。",
                "第三步：生成保守、均衡、进取三种申报策略。",
                "第四步：回测收益、偏差损失和波动。",
                "第五步：模拟储能低价充电、高价放电的增益。",
            ],
        ),
        (
            "5. 交易策略：保守、均衡、进取",
            [
                "保守策略：按预测出力约 85% 申报，优先降低偏差风险。",
                "均衡策略：按预测出力约 92% 申报，兼顾收益和风险。",
                "进取策略：按预测出力约 98% 申报，追求更高收益但承担更高偏差风险。",
            ],
        ),
        (
            "6. 储能增益：不是口号，而是量化测算",
            [
                "午间低价或弃光风险时段充电。",
                "傍晚高价时段放电。",
                "考虑容量、功率和充放电效率约束。",
                "输出有储能与无储能的收益对比，帮助判断储能配置价值。",
            ],
        ),
        (
            "7. 风险边界：先模拟盘，再真实业务",
            [
                "当前使用模拟数据，不代表真实收益。",
                "不同省份偏差考核和交易规则不同，真实落地必须细化规则。",
                "系统建议必须由交易负责人确认，不直接自动下单。",
                "真实部署应先接入只读数据，做模拟盘校准。",
            ],
        ),
        (
            "8. 对隆基的价值：从光伏资产走向收益管理",
            [
                "光伏资产进入电力市场后，交易能力会影响收益质量。",
                "系统可以服务电站运营、现货交易、储能配置和绿电客户服务。",
                "候选人通过这个作品证明：理解业务问题，也能把问题做成工具。",
            ],
        ),
        (
            "9. 候选人表达：我能快速进入真实业务",
            [
                "我不回避缺少实盘经验，所以先做决策辅助和风险复盘。",
                "我能把复杂交易流程拆成数据、模型、策略和复盘模块。",
                "如果进入业务，我会从真实数据接入、模拟盘验证和省份规则细化开始。",
            ],
        ),
    ]
    for title, bullets in slides:
        _add_slide(prs, title, bullets)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"已生成: {OUTPUT}")


if __name__ == "__main__":
    build_presentation()

